# Converts PDFs, Office files, and existing Markdown into cleaned .md files for RAG.
# The script extracts raw text/tables first, then removes OCR/footer/image noise and
# rewrites tables into row-level context so embeddings keep useful document meaning.
# Run with --input and --output; add --overwrite to regenerate existing .md files.

import argparse
import os

# Set before importing OCR/image-heavy libraries.
os.environ.setdefault("OMP_THREAD_LIMIT", "1")

import re
import shutil
import subprocess
import sys
import tempfile
import unicodedata
from collections import Counter
from concurrent.futures import ProcessPoolExecutor, as_completed
from pathlib import Path

MISSING_IMPORTS = {}

try:
    import fitz
except ModuleNotFoundError as exc:
    fitz = None
    MISSING_IMPORTS["fitz"] = exc.name

try:
    import mammoth
except ModuleNotFoundError as exc:
    mammoth = None
    MISSING_IMPORTS["mammoth"] = exc.name

try:
    import openpyxl
except ModuleNotFoundError as exc:
    openpyxl = None
    MISSING_IMPORTS["openpyxl"] = exc.name

try:
    import pytesseract
except ModuleNotFoundError as exc:
    pytesseract = None
    MISSING_IMPORTS["pytesseract"] = exc.name

try:
    from pdf2image import convert_from_path
except ModuleNotFoundError as exc:
    convert_from_path = None
    MISSING_IMPORTS["pdf2image"] = exc.name

try:
    from pptx import Presentation
except ModuleNotFoundError as exc:
    Presentation = None
    MISSING_IMPORTS["python-pptx"] = exc.name

try:
    from tqdm import tqdm
except ModuleNotFoundError:
    def tqdm(iterable=None, total=None, unit=None):
        return iterable if iterable is not None else _NoopProgress(total=total, unit=unit)


class _NoopProgress:
    def __init__(self, total=None, unit=None):
        self.total = total
        self.unit = unit

    def __enter__(self):
        return self

    def __exit__(self, exc_type, exc, tb):
        return False

    def update(self, value):
        return None

    def set_postfix(self, **kwargs):
        return None


EXTENSIONS = {".pdf", ".docx", ".doc", ".rtf", ".xlsx", ".pptx", ".md"}
DEFAULT_MAX_WORKERS = max(1, min((os.cpu_count() or 4) - 1, 8))
OCR_DPI = 300
OCR_TIMEOUT_S = 300
TABLE_ROW_CHUNK_SIZE = 10
MAX_EXCEL_HEADER_SCAN_ROWS = 8


BASE64_RE = re.compile(r"!\[[^\]]*\]\(data:image/[^)]+\)")
PICTURE_RE = re.compile(r"[*_]{0,2}\s*==>\s*[^<\n]*?<==\s*[*_]{0,2}")
PAGE_NUM_RE = re.compile(r"^\s*\d{1,4}\s*$", re.MULTILINE)
PAGE_NUM_DASH_RE = re.compile(r"^\s*[-—–~_=]+\s*\d{1,4}\s*[-—–~_=]+\s*$", re.MULTILINE)
PAGE_FRAC_RE = re.compile(r"^\s*\d{1,4}\s*/\s*\d{1,4}\s*$", re.MULTILINE)
HYPHEN_BREAK_RE = re.compile(r"(\w+)-\n(\w+)")
MULTI_BLANK_RE = re.compile(r"\n{3,}")
BR_TAG_RE = re.compile(r"<\s*br\s*/?\s*>", re.IGNORECASE)
HEADING_RE = re.compile(r"^(#{1,6})\s+(.+?)\s*$")
NUMERIC_ONLY_RE = re.compile(r"^[\s\d.,:/\\|+%°№#()=-]+$")
GRAPH_WORD_RE = re.compile(
    r"\b(?:скважина|керн|кп|проницаем|непроницаем|граничн|значени|график|диаграмм|chart|figure)\b",
    re.IGNORECASE,
)

PICTURE_TEXT_PAIR_RE = re.compile(
    r"\*{0,2}\s*-{3,}\s*(?:Beginning|Start)\s+of\s+picture\s+text\s*-{3,}\s*\*{0,2}"
    r".*?"
    r"\*{0,2}\s*-{3,}\s*End\s+of\s+picture\s+text\s*-{3,}\s*\*{0,2}",
    re.IGNORECASE | re.DOTALL,
)
PICTURE_TEXT_END_RE = re.compile(
    r"\*{0,2}\s*-{3,}\s*End\s+of\s+picture\s+text\s*-{3,}\s*\*{0,2}",
    re.IGNORECASE,
)
TABLE_SEP_RE = re.compile(r"^\|?\s*[:\-]+\s*(\|\s*[:\-]+\s*)*\|?\s*$")
FIGURE_CAPTION_RE = re.compile(
    r"^\*{0,2}\s*(?:Рис(?:унок|\.)?|Fig(?:ure|\.)?)\s*[\d.\-]*\s*[—\-–]?\s*[^\n]{0,200}\*{0,2}$",
    re.IGNORECASE,
)
TABLE_CAPTION_RE = re.compile(
    r"^\*{0,2}\s*(?:Табл(?:ица|\.)?|Table)\s*[\d.\-]*\s*[—\-–:.]?\s*[^\n]{0,250}\*{0,2}$",
    re.IGNORECASE,
)


def sanitize(text: str) -> str:
    text = unicodedata.normalize("NFC", text)
    text = (
        text.replace("ﬁ", "fi")
        .replace("ﬂ", "fl")
        .replace("ﬀ", "ff")
        .replace("ﬃ", "ffi")
        .replace("ﬄ", "ffl")
    )
    text = (
        text.replace("\u00ad", "")
        .replace("\u200b", "")
        .replace("\ufeff", "")
        .replace("\u00a0", " ")
    )
    return text.encode("utf-8", "ignore").decode("utf-8")


def is_meaningful_text(text: str, min_chars: int = 80, min_alpha_ratio: float = 0.2) -> bool:
    text = text.strip()
    if len(text) < min_chars:
        return False
    alpha = sum(1 for char in text if char.isalpha())
    return (alpha / max(len(text), 1)) >= min_alpha_ratio


def pdf_has_good_text(pdf_path: Path) -> bool:
    with fitz.open(str(pdf_path)) as doc:
        if len(doc) == 0:
            return False
        page_texts = [page.get_text() for page in doc]
        joined = "\n".join(page_texts)
        if is_meaningful_text(joined, min_chars=120, min_alpha_ratio=0.18):
            return True
        good_pages = sum(1 for text in page_texts if is_meaningful_text(text))
        return good_pages / len(doc) >= 0.35


def pdf_embedded_text_to_markdown(pdf_path: Path) -> str:
    parts = [f"# {pdf_path.stem}"]
    with fitz.open(str(pdf_path)) as doc:
        for page_num, page in enumerate(doc, 1):
            page_text = page.get_text("text").strip()
            if page_text:
                parts.append(f"## Страница {page_num}\n\n{page_text}")
    return "\n\n".join(parts)


def escape_cell(value) -> str:
    if value is None:
        return ""
    return str(value).replace("|", "\\|").replace("\n", " ").strip()


def clean_cell(value: str) -> str:
    value = BR_TAG_RE.sub(" ", value)
    value = re.sub(r"\*\*(.+?)\*\*", r"\1", value)
    value = re.sub(r"__(.+?)__", r"\1", value)
    value = re.sub(r"(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)", r"\1", value)
    value = re.sub(r"_(.+?)_", r"\1", value)
    value = re.sub(r"\s+", " ", value)
    return value.strip()


def strip_outer_table_pipes(line: str) -> str:
    stripped = line.strip()
    if stripped.startswith("|"):
        stripped = stripped[1:]
    if stripped.endswith("|") and not stripped.endswith("\\|"):
        stripped = stripped[:-1]
    return stripped


def split_markdown_row(line: str) -> list[str]:
    cells = []
    current = []
    text = strip_outer_table_pipes(line)
    idx = 0
    while idx < len(text):
        char = text[idx]
        if char == "\\" and idx + 1 < len(text) and text[idx + 1] == "|":
            current.append("|")
            idx += 2
            continue
        if char == "|":
            cells.append("".join(current))
            current = []
            idx += 1
            continue
        current.append(char)
        idx += 1
    cells.append("".join(current))
    return [clean_cell(cell) for cell in cells]


def context_prefix(doc_name: str, section: str | None = None, caption: str | None = None) -> str:
    parts = [f"Документ: {doc_name}"]
    if section:
        parts.append(f"Раздел: {section}")
    if caption:
        parts.append(caption.strip().strip("*").strip())
    return "[" + " | ".join(parts) + "]"


def chunk_rows(rows: list[str], prefix: str) -> str:
    chunks = []
    for index in range(0, len(rows), TABLE_ROW_CHUNK_SIZE):
        chunk = rows[index : index + TABLE_ROW_CHUNK_SIZE]
        chunks.append(prefix + "\n" + "\n".join(chunk))
    return "\n\n".join(chunks)


# Turns table rows into compact "header: value" text chunks for better RAG retrieval.
def rows_to_rag_text(headers: list[str], data_rows: list[list[str]], prefix: str) -> str:
    headers = [clean_cell(str(header)) for header in headers]
    max_width = max([len(headers), *[len(row) for row in data_rows]], default=0)
    if len(headers) < max_width:
        headers.extend([f"Колонка {idx + 1}" for idx in range(len(headers), max_width)])

    out = []
    for row in data_rows:
        cells = [clean_cell(str(cell)) for cell in row]
        if len(cells) < len(headers):
            cells.extend([""] * (len(headers) - len(cells)))
        if not any(cells):
            continue

        non_empty = [(idx, cell) for idx, cell in enumerate(cells) if cell and cell.strip("-—–")]
        if len(non_empty) == 1 and any(char.isalpha() for char in non_empty[0][1]):
            out.append(f"[{non_empty[0][1]}]")
            continue

        pairs = [
            f"{header}: {cell}"
            for header, cell in zip(headers, cells)
            if header and cell and cell.strip("-—–")
        ]
        if pairs:
            out.append("; ".join(pairs))
            continue

        values = [cell for cell in cells if cell and cell.strip("-—–")]
        if values:
            out.append(" — ".join(values))

    return chunk_rows(out, prefix) if out else ""


def markdown_table_to_rag_text(block: str, prefix: str) -> str:
    lines = [line for line in block.splitlines() if line.strip().startswith("|")]
    if len(lines) < 2:
        return ""

    sep_idx = None
    for idx, line in enumerate(lines):
        if TABLE_SEP_RE.match(line.strip()):
            sep_idx = idx
            break

    if sep_idx is not None and sep_idx > 0:
        header_rows = [split_markdown_row(line) for line in lines[:sep_idx]]
        column_count = max(len(row) for row in header_rows)
        headers = []
        for col_idx in range(column_count):
            parts = []
            for row in header_rows:
                if col_idx < len(row) and row[col_idx]:
                    parts.append(row[col_idx])
            deduped = []
            for part in parts:
                if not deduped or deduped[-1] != part:
                    deduped.append(part)
            headers.append(" ".join(deduped))
        data_lines = lines[sep_idx + 1 :]
    else:
        headers = split_markdown_row(lines[0])
        data_lines = lines[1:]

    data_rows = [
        split_markdown_row(line)
        for line in data_lines
        if not TABLE_SEP_RE.match(line.strip())
    ]
    return rows_to_rag_text(headers, data_rows, prefix)


def row_signal(row: list[str]) -> tuple[int, int, int]:
    non_empty = [cell for cell in row if clean_cell(str(cell))]
    alpha_cells = [cell for cell in non_empty if any(char.isalpha() for char in str(cell))]
    numeric_cells = [cell for cell in non_empty if any(char.isdigit() for char in str(cell))]
    return len(non_empty), len(alpha_cells), len(numeric_cells)


def choose_excel_header_index(rows: list[list[str]]) -> int:
    scan_limit = min(len(rows), MAX_EXCEL_HEADER_SCAN_ROWS)
    best_index = 0
    best_score = (-1, -1, -1)
    for idx in range(scan_limit):
        row = rows[idx]
        non_empty, alpha_count, numeric_count = row_signal(row)
        if non_empty < 2:
            continue
        next_non_empty = 0
        if idx + 1 < len(rows):
            next_non_empty = row_signal(rows[idx + 1])[0]
        score = (alpha_count, min(non_empty, next_non_empty), -numeric_count)
        if score > best_score:
            best_index = idx
            best_score = score
    return best_index


def merge_excel_header_rows(header_rows: list[list[str]]) -> list[str]:
    width = max((len(row) for row in header_rows), default=0)
    headers = []
    for col_idx in range(width):
        parts = []
        for row in header_rows:
            if col_idx < len(row):
                cell = clean_cell(str(row[col_idx]))
                if cell and (not parts or parts[-1] != cell):
                    parts.append(cell)
        headers.append(" ".join(parts).strip() or f"Колонка {col_idx + 1}")
    return headers


def build_excel_headers(rows: list[list[str]], header_idx: int) -> list[str]:
    start_idx = header_idx
    for idx in range(header_idx - 1, -1, -1):
        non_empty, alpha_count, _ = row_signal(rows[idx])
        next_non_empty = row_signal(rows[idx + 1])[0]
        if non_empty >= 2 and alpha_count >= 2 and non_empty >= max(2, next_non_empty // 2):
            start_idx = idx
            continue
        break
    return merge_excel_header_rows(rows[start_idx : header_idx + 1])


def is_pipe_table_paragraph(paragraph: str) -> bool:
    lines = [line for line in paragraph.splitlines() if line.strip()]
    if len(lines) < 2:
        return False
    return all(line.lstrip().startswith("|") for line in lines)


def find_recurring_lines(text: str) -> set[str]:
    lines = [line.strip() for line in text.splitlines()]
    candidates = [
        line
        for line in lines
        if 8 <= len(line) <= 200
        and not line.startswith(("#", "|", "```", ">"))
        and not re.fullmatch(r"[\d\W_]+", line)
    ]
    counts = Counter(candidates)
    return {line for line, count in counts.items() if count >= 5}


def paragraph_has_value(paragraph: str) -> bool:
    stripped = paragraph.strip()
    if not stripped:
        return False
    if HEADING_RE.match(stripped):
        return any(char.isalpha() for char in stripped)
    if stripped.startswith("[Документ:"):
        return True
    if re.match(r"^[-*+]\s+\S+", stripped):
        return any(char.isalpha() for char in stripped)
    if re.match(r"^[A-Za-zА-Яа-яЁё0-9][\w .:/\\№#-]{0,80}:\s*\S+", stripped):
        return True
    if re.fullmatch(r"[\w .:/\\№#-]{2,120}", stripped) and any(char.isdigit() for char in stripped):
        return True
    return sum(1 for char in stripped if char.isalpha()) >= 10


def drop_picture_text_blocks(text: str, stats: dict) -> str:
    text, dropped = PICTURE_TEXT_PAIR_RE.subn("", text)
    stats["picture_blocks_dropped"] += dropped

    kept = []
    for paragraph in re.split(r"\n{2,}", text):
        if PICTURE_TEXT_END_RE.search(paragraph):
            stats["picture_blocks_dropped"] += 1
            continue
        br_count = len(BR_TAG_RE.findall(paragraph))
        sentence_count = len(re.findall(r"[.!?]\s+[А-ЯA-Z]", paragraph))
        if br_count >= 8 and sentence_count < 2:
            stats["picture_blocks_dropped"] += 1
            continue
        kept.append(paragraph)
    return "\n\n".join(kept)


def looks_like_ocr_graph_block(paragraph: str) -> bool:
    lines = [line.strip() for line in paragraph.splitlines() if line.strip()]
    if len(lines) < 8:
        return False

    numeric_lines = [line for line in lines if NUMERIC_ONLY_RE.fullmatch(line)]
    short_lines = [line for line in lines if len(line) <= 30]
    graph_word_hits = sum(1 for line in lines if GRAPH_WORD_RE.search(line))

    numeric_ratio = len(numeric_lines) / len(lines)
    short_ratio = len(short_lines) / len(lines)

    if graph_word_hits >= 2 and numeric_ratio >= 0.35 and short_ratio >= 0.65:
        return True
    if len(numeric_lines) >= 6 and short_ratio >= 0.75 and graph_word_hits >= 1:
        return True
    return False


def drop_ocr_graph_blocks(text: str, stats: dict) -> str:
    kept = []
    for paragraph in re.split(r"\n{2,}", text):
        if looks_like_ocr_graph_block(paragraph):
            stats["graph_blocks_dropped"] += 1
            continue
        kept.append(paragraph)
    return "\n\n".join(kept)


def process_tables_with_context(text: str, doc_name: str, stats: dict) -> str:
    paragraphs = re.split(r"\n{2,}", text)
    current_section = None
    pending_caption = None
    out = []

    for paragraph in paragraphs:
        stripped = paragraph.strip()
        if not stripped:
            continue

        heading = HEADING_RE.match(stripped)
        if heading:
            current_section = heading.group(2).strip()
            if pending_caption:
                out.append(pending_caption)
                pending_caption = None
            out.append(paragraph)
            continue

        if (
            TABLE_CAPTION_RE.fullmatch(stripped)
            and not FIGURE_CAPTION_RE.fullmatch(stripped)
            and len(stripped) < 300
        ):
            if pending_caption:
                out.append(pending_caption)
            pending_caption = stripped
            continue

        if is_pipe_table_paragraph(paragraph):
            prefix = context_prefix(doc_name, current_section, pending_caption)
            pending_caption = None
            converted = markdown_table_to_rag_text(paragraph, prefix)
            if converted:
                stats["tables_converted"] += 1
                out.append(converted)
            continue

        if pending_caption:
            out.append(pending_caption)
            pending_caption = None
        out.append(paragraph)

    if pending_caption:
        out.append(pending_caption)
    return "\n\n".join(out)


# Main cleanup pass before writing Markdown: normalize text, remove noise, convert
# tables, and keep only paragraphs that are likely to help search/answers.
def clean_markdown(text: str, doc_name: str) -> tuple[str, dict]:
    stats = {
        "footers_removed": 0,
        "paras_dropped": 0,
        "tables_converted": 0,
        "picture_blocks_dropped": 0,
        "graph_blocks_dropped": 0,
        "figure_captions_dropped": 0,
    }

    text = sanitize(text)
    text = BASE64_RE.sub("", text)
    text = PICTURE_RE.sub("", text)
    text = drop_picture_text_blocks(text, stats)
    text = drop_ocr_graph_blocks(text, stats)

    paragraphs = []
    for paragraph in re.split(r"\n{2,}", text):
        stripped = paragraph.strip()
        if FIGURE_CAPTION_RE.fullmatch(stripped) and len(stripped) < 200:
            stats["figure_captions_dropped"] += 1
            continue
        paragraphs.append(paragraph)
    text = "\n\n".join(paragraphs)

    text = process_tables_with_context(text, doc_name, stats)
    text = HYPHEN_BREAK_RE.sub(r"\1\2", text)
    text = PAGE_NUM_RE.sub("", text)
    text = PAGE_NUM_DASH_RE.sub("", text)
    text = PAGE_FRAC_RE.sub("", text)

    footers = find_recurring_lines(text)
    if footers:
        kept = []
        for line in text.splitlines():
            if line.strip() in footers:
                stats["footers_removed"] += 1
                continue
            kept.append(line)
        text = "\n".join(kept)

    good = []
    for paragraph in re.split(r"\n{2,}", text):
        if paragraph_has_value(paragraph):
            good.append(paragraph.strip())
        else:
            stats["paras_dropped"] += 1

    text = "\n\n".join(good)
    text = MULTI_BLANK_RE.sub("\n\n", text)
    return text.strip() + "\n", stats


# Prefer embedded PDF text; OCR each page only when the PDF looks scanned and OCR is allowed.
def pdf_to_raw_markdown(pdf_path: Path, no_ocr: bool, poppler_path: str | None) -> tuple[str, str]:
    if pdf_has_good_text(pdf_path):
        markdown = pdf_embedded_text_to_markdown(pdf_path)
        if is_meaningful_text(markdown):
            return markdown, "pdf"
        if no_ocr:
            return markdown, "pdf-low-text"

    if no_ocr:
        raise ValueError("PDF looks scanned or empty; OCR is disabled")

    with fitz.open(str(pdf_path)) as doc:
        page_count = len(doc)

    parts = [f"# {pdf_path.stem}"]
    for page_num in range(1, page_count + 1):
        images = convert_from_path(
            str(pdf_path),
            dpi=OCR_DPI,
            poppler_path=poppler_path,
            first_page=page_num,
            last_page=page_num,
        )
        page_text = pytesseract.image_to_string(images[0], lang="rus+eng", timeout=OCR_TIMEOUT_S)
        if page_text.strip():
            parts.append(f"## Страница {page_num}\n\n{page_text.strip()}")
        del images
    return "\n\n".join(parts), "pdf-ocr"


def docx_to_raw_markdown(docx_path: Path) -> str:
    with open(docx_path, "rb") as file:
        result = mammoth.convert_to_markdown(file)
    return result.value


def doc_to_raw_markdown(doc_path: Path, libreoffice: str) -> str:
    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        profile_uri = (tmp_path / "lo_profile").as_uri()
        result = subprocess.run(
            [
                libreoffice,
                f"-env:UserInstallation={profile_uri}",
                "--headless",
                "--convert-to",
                "docx",
                "--outdir",
                str(tmp_path),
                str(doc_path),
            ],
            capture_output=True,
            timeout=180,
        )
        if result.returncode != 0:
            err = result.stderr.decode("utf-8", errors="ignore")
            raise RuntimeError(f"LibreOffice code {result.returncode}: {err}")

        produced = list(tmp_path.glob("*.docx"))
        if not produced:
            raise RuntimeError("LibreOffice did not create a .docx file")
        return docx_to_raw_markdown(produced[0])


def xlsx_to_raw_markdown(xlsx_path: Path) -> str:
    wb = openpyxl.load_workbook(str(xlsx_path), data_only=True, read_only=True)
    parts = [f"# {xlsx_path.stem}"]
    try:
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            rows = [
                [escape_cell(cell) for cell in row]
                for row in ws.iter_rows(values_only=True)
                if any(cell is not None for cell in row)
            ]
            if not rows:
                continue

            header_idx = choose_excel_header_index(rows)
            headers = build_excel_headers(rows, header_idx)
            data_rows = rows[header_idx + 1 :]
            prefix = context_prefix(xlsx_path.stem, f"Лист: {sheet}", "Таблица из Excel")
            table_text = rows_to_rag_text(headers, data_rows, prefix)
            if table_text:
                parts.append(f"## {sheet}\n\n{table_text}")
    finally:
        wb.close()
    return "\n\n".join(parts)


def iter_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == 6:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def pptx_to_raw_markdown(pptx_path: Path) -> str:
    presentation = Presentation(str(pptx_path))
    parts = [f"# {pptx_path.stem}"]
    for slide_num, slide in enumerate(presentation.slides, 1):
        slide_parts = [f"## Слайд {slide_num}"]
        table_num = 0
        for shape in iter_shapes(slide.shapes):
            if shape.has_text_frame and shape.text_frame.text.strip():
                slide_parts.append(shape.text_frame.text.strip())
            elif getattr(shape, "has_table", False) and shape.has_table:
                table_num += 1
                rows = [
                    [escape_cell(cell.text) for cell in row.cells]
                    for row in shape.table.rows
                ]
                if rows:
                    prefix = context_prefix(
                        pptx_path.stem,
                        f"Слайд {slide_num}",
                        f"Таблица {table_num} из презентации",
                    )
                    table_text = rows_to_rag_text(rows[0], rows[1:], prefix)
                    if table_text:
                        slide_parts.append(table_text)

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_parts.append(f"Заметки: {notes}")

        if len(slide_parts) > 1:
            parts.append("\n\n".join(slide_parts))
    return "\n\n".join(parts)


# Sends each supported file type to the extractor that understands it.
def convert_to_raw_markdown(src: Path, config: dict) -> tuple[str, str]:
    ext = src.suffix.lower()
    if ext == ".pdf":
        return pdf_to_raw_markdown(src, config["no_ocr"], config["poppler_path"])
    if ext == ".docx":
        return docx_to_raw_markdown(src), "docx"
    if ext in {".doc", ".rtf"}:
        return doc_to_raw_markdown(src, config["libreoffice"]), ext[1:]
    if ext == ".xlsx":
        return xlsx_to_raw_markdown(src), "xlsx"
    if ext == ".pptx":
        return pptx_to_raw_markdown(src), "pptx"
    if ext == ".md":
        return src.read_text(encoding="utf-8", errors="ignore"), "md-clean"
    raise ValueError(f"Unsupported extension: {ext}")


# One source file becomes one cleaned .md file, preserving the input folder structure.
def process_file(src_str: str, input_dir: str, output_dir: str, config: dict) -> tuple[str, str, str, dict]:
    src = Path(src_str)
    input_path = Path(input_dir)
    output_path = Path(output_dir)
    empty_stats = {
        "footers_removed": 0,
        "paras_dropped": 0,
        "tables_converted": 0,
        "picture_blocks_dropped": 0,
        "graph_blocks_dropped": 0,
        "figure_captions_dropped": 0,
    }

    try:
        rel = src.relative_to(input_path)
        out_dir = output_path / rel.parent
        out_dir.mkdir(parents=True, exist_ok=True)
        out_file = out_dir / f"{src.stem}.md"

        if out_file.exists() and out_file.stat().st_size > 0 and not config["overwrite"]:
            return str(src), "skip", "exists", empty_stats

        raw_markdown, mode = convert_to_raw_markdown(src, config)
        cleaned, stats = clean_markdown(raw_markdown, src.stem)
        if not cleaned.strip():
            raise ValueError("Resulting markdown is empty after cleaning")

        out_file.write_text(cleaned, encoding="utf-8")
        return str(src), "ok", mode, stats
    except Exception as exc:
        return str(src), "fail", f"{type(exc).__name__}: {exc}", empty_stats


def resolve_executable(value: str | None, names: list[str]) -> str | None:
    if value:
        return value
    for name in names:
        found = shutil.which(name)
        if found:
            return found
    return None


def validate_args(args, files: list[Path]) -> dict:
    config = {
        "overwrite": args.overwrite,
        "no_ocr": args.no_ocr,
        "poppler_path": args.poppler_path,
        "libreoffice": resolve_executable(
            args.libreoffice,
            ["soffice", "libreoffice", "soffice.exe"],
        ),
    }

    if args.tesseract_cmd and pytesseract is not None:
        pytesseract.pytesseract.tesseract_cmd = args.tesseract_cmd

    extensions = {file.suffix.lower() for file in files}
    errors = []

    missing_by_extension = {
        ".pdf": ["fitz"],
        ".docx": ["mammoth"],
        ".doc": ["mammoth"],
        ".rtf": ["mammoth"],
        ".xlsx": ["openpyxl"],
        ".pptx": ["python-pptx"],
    }
    for extension in extensions:
        for package in missing_by_extension.get(extension, []):
            if package in MISSING_IMPORTS:
                errors.append(f"Python package is missing for {extension}: {package}")

    if ".pdf" in extensions and not args.no_ocr:
        for package in ["pytesseract", "pdf2image"]:
            if package in MISSING_IMPORTS:
                errors.append(f"Python package is missing for PDF OCR: {package}")

    if {".doc", ".rtf"} & extensions and not config["libreoffice"]:
        errors.append("LibreOffice is required for .doc/.rtf. Pass --libreoffice or add soffice to PATH.")

    if ".pdf" in extensions and not args.no_ocr and pytesseract is not None:
        tesseract = pytesseract.pytesseract.tesseract_cmd
        if tesseract != "tesseract" and not Path(tesseract).exists():
            errors.append(f"Tesseract not found: {tesseract}")
        elif tesseract == "tesseract" and not shutil.which("tesseract"):
            errors.append("Tesseract not found in PATH. Pass --tesseract-cmd or use --no-ocr.")

    if errors:
        for error in errors:
            print(f"ERROR: {error}", file=sys.stderr)
        sys.exit(1)

    return config


def parse_args():
    parser = argparse.ArgumentParser(
        description="Convert documents to clean Markdown for RAG/embeddings."
    )
    parser.add_argument("--input", "-i", required=True, help="Input folder with source files.")
    parser.add_argument("--output", "-o", required=True, help="Output folder for .md files.")
    parser.add_argument("--failed-log", help="Path for failed conversions log.")
    parser.add_argument("--workers", "-w", type=int, default=DEFAULT_MAX_WORKERS)
    parser.add_argument("--overwrite", action="store_true", help="Overwrite existing .md files.")
    parser.add_argument("--no-ocr", action="store_true", help="Disable OCR fallback for scanned PDFs.")
    parser.add_argument("--tesseract-cmd", help="Path to tesseract executable.")
    parser.add_argument("--poppler-path", help="Path to Poppler bin folder for pdf2image.")
    parser.add_argument("--libreoffice", help="Path to LibreOffice soffice executable.")
    return parser.parse_args()


# Command-line workflow: find files, validate external tools/packages, process files,
# then print a short conversion report.
def main():
    args = parse_args()
    input_dir = Path(args.input).expanduser().resolve()
    output_dir = Path(args.output).expanduser().resolve()
    failed_log = Path(args.failed_log).expanduser().resolve() if args.failed_log else output_dir / "failed_convert.txt"

    if not input_dir.exists():
        print(f"ERROR: input folder does not exist: {input_dir}", file=sys.stderr)
        sys.exit(1)

    files = [
        file
        for file in input_dir.rglob("*")
        if file.is_file() and file.suffix.lower() in EXTENSIONS
    ]
    config = validate_args(args, files)
    output_dir.mkdir(parents=True, exist_ok=True)

    print(f"Files: {len(files)} | workers: {args.workers}")
    if not files:
        print(f"Output: {output_dir}")
        print("No supported files found.")
        sys.exit(0)

    ok = skip = fail = 0
    failed = []
    totals = Counter()

    def handle_result(src: str, status: str, message: str, stats: dict):
        nonlocal ok, skip, fail
        if status == "ok":
            ok += 1
            totals.update(stats)
        elif status == "skip":
            skip += 1
        else:
            fail += 1
            failed.append(f"{src}\t{message}")

    def run_sequential():
        with tqdm(total=len(files), unit="file") as progress:
            for file in files:
                src, status, message, stats = process_file(str(file), str(input_dir), str(output_dir), config)
                handle_result(src, status, message, stats)
                progress.set_postfix(ok=ok, skip=skip, fail=fail)
                progress.update(1)

    def run_parallel(worker_count: int):
        with ProcessPoolExecutor(max_workers=worker_count) as executor:
            futures = {
                executor.submit(process_file, str(file), str(input_dir), str(output_dir), config): file
                for file in files
            }
            with tqdm(total=len(files), unit="file") as progress:
                for future in as_completed(futures):
                    src, status, message, stats = future.result()
                    handle_result(src, status, message, stats)
                    progress.set_postfix(ok=ok, skip=skip, fail=fail)
                    progress.update(1)

    worker_count = max(1, args.workers)
    if worker_count == 1:
        run_sequential()
    else:
        try:
            run_parallel(worker_count)
        except (OSError, PermissionError) as exc:
            print(f"Parallel mode failed ({type(exc).__name__}: {exc}). Falling back to one worker.")
            ok = skip = fail = 0
            failed.clear()
            totals.clear()
            run_sequential()

    print(f"\nOK: {ok} | skipped: {skip} | failed: {fail}")
    print(f"Output: {output_dir}")
    print(f"Tables converted to RAG text: {totals['tables_converted']}")
    print(f"Footers removed: {totals['footers_removed']}")
    print(f"Low-value paragraphs dropped: {totals['paras_dropped']}")
    print(f"Picture OCR blocks dropped: {totals['picture_blocks_dropped']}")
    print(f"Graph OCR blocks dropped: {totals['graph_blocks_dropped']}")
    print(f"Figure captions dropped: {totals['figure_captions_dropped']}")

    if failed:
        failed_log.write_text("\n".join(failed), encoding="utf-8")
        print(f"Failed log: {failed_log}")

    sys.exit(0 if fail == 0 else 2)


if __name__ == "__main__":
    main()
